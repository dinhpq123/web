#!/usr/bin/env python3
"""
Quawaco IOC Center – Professional PPTX Generator
Generates docs/Quawaco_IOC_Presentation.pptx
"""
import os, sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")

# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent
IMGS = SCRIPT_DIR / "screenshots"
OUT  = SCRIPT_DIR / "Quawaco_IOC_Presentation.pptx"

# ── Color Palette ──────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x03, 0x0e, 0x1c)   # deep navy
C_BLUE   = RGBColor(0x00, 0x50, 0xcc)   # primary blue
C_CYAN   = RGBColor(0x00, 0xc8, 0xff)   # accent cyan
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_MUTED  = RGBColor(0x54, 0x6e, 0x7a)
C_GREEN  = RGBColor(0x00, 0xe6, 0x76)
C_YELLOW = RGBColor(0xff, 0xca, 0x28)
C_RED    = RGBColor(0xff, 0x17, 0x44)
C_GRAY   = RGBColor(0x1e, 0x2d, 0x3d)

# ── Slide size: 16:9 widescreen ────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def blank(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=C_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color
    return txb

def add_img(slide, path, l, t, w, h=None):
    p = IMGS / path if not Path(path).is_absolute() else Path(path)
    # find with timestamp suffix if exact match fails
    if not p.exists():
        matches = list(IMGS.glob(f"{p.stem}_*.png"))
        if matches:
            p = sorted(matches)[-1]
    if not p.exists():
        print(f"WARNING: Image not found: {p}")
        return None
    if h:
        return slide.shapes.add_picture(str(p), l, t, w, h)
    return slide.shapes.add_picture(str(p), l, t, w)

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, W, Inches(1.3), C_GRAY)
    add_rect(slide, 0, Inches(1.28), W, Emu(60000), C_CYAN)
    add_text(slide, title, Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.65),
             size=26, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72), W - Inches(0.8), Inches(0.5),
                 size=13, color=C_CYAN)

def footer(slide, text="Quawaco IOC Center  |  Tài liệu mật – Chỉ dùng nội bộ  |  © 2026"):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), C_GRAY)
    add_text(slide, text, Inches(0.3), H - Inches(0.32), W - Inches(0.6), Inches(0.3),
             size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

def screen_slide(prs, title, subtitle, img_filename, caption):
    sl = blank(prs)
    bg(sl)
    header_bar(sl, title, subtitle)
    img = add_img(sl, img_filename, Inches(0.3), Inches(1.42), Inches(12.7), Inches(5.5))
    add_rect(sl, 0, H - Inches(0.7), W, Inches(0.35), C_GRAY)
    add_text(sl, f"📸  {caption}", Inches(0.3), H - Inches(0.68), W - Inches(0.6), Inches(0.32),
             size=9, color=C_CYAN, align=PP_ALIGN.LEFT, italic=True)
    footer(sl)
    return sl

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – COVER
# ══════════════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl)
add_rect(sl, 0, 0, W, H, C_DARK)
add_rect(sl, 0, 0, Inches(0.18), H, C_CYAN)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.4), C_GRAY)
add_rect(sl, Inches(0.4), Inches(1.2), Inches(12.5), Emu(40000), C_CYAN)

add_text(sl, "QUAWACO IOC CENTER", Inches(0.8), Inches(1.5), Inches(12), Inches(1),
         size=42, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Tài liệu Đặc tả Tính năng & Thuyết trình Hệ thống",
         Inches(0.8), Inches(2.55), Inches(12), Inches(0.7),
         size=20, color=C_CYAN, align=PP_ALIGN.LEFT)
footer(sl)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – MỤC LỤC
# ══════════════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl)
header_bar(sl, "Mục lục hệ thống", "6 Phân hệ Điều hành chính")
modules = [
    ("01", "Điều hành Tập trung (Command Center)", C_CYAN),
    ("02", "Quản lý Sản xuất (Production MGMT)", C_BLUE),
    ("03", "Kỹ thuật & Mạng lưới (Engineering)", C_GREEN),
    ("04", "Quản lý Kinh doanh (Business MGMT)", C_YELLOW),
    ("05", "Nguồn lực & Hệ thống (Resources)", C_RED),
    ("06", "Trung tâm AI (AI Center)", C_CYAN),
]
for idx, (num, name, col) in enumerate(modules):
    lft = Inches(1.5)
    tp  = Inches(1.8) + idx * Inches(0.8)
    add_rect(sl, lft, tp, Inches(0.6), Inches(0.6), col)
    add_text(sl, num, lft, tp + Inches(0.1), Inches(0.6), Inches(0.4),
             size=18, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(sl, name, lft + Inches(0.8), tp + Inches(0.1), Inches(8), Inches(0.6),
             size=18, color=C_WHITE)
footer(sl)

# ══════════════════════════════════════════════════════════════════════════
# MODULE 1: ĐIỀU HÀNH TẬP TRUNG
# ══════════════════════════════════════════════════════════════════════════
screen_slide(prs, "01 · Điều hành Tập trung – Dashboard",
    "Giao diện giám sát 360 độ · Tích hợp Alarms panel & Ticker",
    "02_dashboard_collapsed.png",
    "Dashboard tổng quan: KPI sản lượng, trạm trực tuyến, NRW và bản đồ tiêu thụ")

screen_slide(prs, "01 · Điều hành Tập trung – Bản đồ GIS",
    "Quản lý hạ tầng trên nền OSM · Lớp đường ống & DMA",
    "04_gis.png",
    "Bản đồ GIS: Hiển thị mạng lưới tuyến ống, trạm bơm và các điểm sự cố thời gian thực")

screen_slide(prs, "01 · Điều hành Tập trung – Camera CCTV",
    "Giám sát an ninh & vận hành qua hình ảnh trực tiếp",
    "05_camera.png",
    "Hệ thống Camera: Tích hợp luồng video từ các nhà máy và trạm bơm trọng điểm")

# ══════════════════════════════════════════════════════════════════════════
# MODULE 2: QUẢN LÝ SẢN XUẤT
# ══════════════════════════════════════════════════════════════════════════
screen_slide(prs, "02 · Quản lý Sản xuất – Giám sát SCADA",
    "Thông số Áp lực · Lưu lượng · Mực nước bể chứa",
    "07_scada.png",
    "SCADA: Giám sát toàn bộ các trạm bơm với biểu đồ xu hướng 24h")

screen_slide(prs, "02 · Quản lý Sản xuất – Lịch Bơm (Scheduler)",
    "Tối ưu hóa giờ chạy bơm · Tiết kiệm điện năng AI",
    "11_scheduler.png",
    "Lịch Bơm: Bảng phân bổ giờ chạy bơm theo khung giá điện EVN, hỗ trợ reset AI")

screen_slide(prs, "02 · Quản lý Sản xuất – Chất lượng nước",
    "Giám sát 6 chỉ tiêu QCVN 01:2009 · Cảnh báo tức thời",
    "09_quality.png",
    "Chất lượng nước: Theo dõi pH, Độ đục, Clo dư... với đồng hồ gauge trực quan")

# ══════════════════════════════════════════════════════════════════════════
# MODULE 3: KỸ THUẬT & MẠNG LƯỚI
# ══════════════════════════════════════════════════════════════════════════
screen_slide(prs, "03 · Kỹ thuật & Mạng lưới – Quản lý Sự cố",
    "Vòng đời xử lý sự cố · Lệnh công tác hiện trường",
    "14_incidents.png",
    "Sự cố: Quản lý tiếp nhận, phân công và đóng sự cố với đầy đủ timeline & photo")

screen_slide(prs, "03 · Kỹ thuật & Mạng lưới – NRW Thất thoát",
    "Phân tích cân bằng nước theo DMA · Cảnh báo rò rỉ",
    "13_nrw.png",
    "Quản lý Thất thoát: Tỷ lệ NRW theo vùng, hỗ trợ chiến dịch giảm thất thoát nước")

# ══════════════════════════════════════════════════════════════════════════
# MODULE 4: QUẢN LÝ KINH DOANH
# ══════════════════════════════════════════════════════════════════════════
screen_slide(prs, "04 · Quản lý Kinh doanh – Khách hàng & HĐ",
    "Quản lý 49k+ khách hàng · Hợp đồng điện tử · Lịch sử ghi thu",
    "17_customers.png",
    "Khách hàng: Tra cứu thông tin hồ sơ khách hàng 360 độ và lịch sử tiêu thụ")

# ══════════════════════════════════════════════════════════════════════════
# MODULE 6: TRUNG TÂM AI
# ══════════════════════════════════════════════════════════════════════════
screen_slide(prs, "06 · Trung tâm AI – AI Agent",
    "Phân tích dữ liệu bằng ngôn ngữ tự nhiên · Dự báo xu hướng",
    "22_aiagent.png",
    "AI Agent: Trợ lý phân tích dữ liệu thông minh, phát hiện bất thường tự động")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE: THANK YOU
# ══════════════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl)
add_text(sl, "Cảm ơn Quý đối tác!", Inches(0.8), Inches(3.0), Inches(12), Inches(1.0),
         size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
footer(sl)

# ── Save ───────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"DONE  {OUT}  ({OUT.stat().st_size//1024} KB, {len(prs.slides)} slides)")
